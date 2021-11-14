from pptx import Presentation
from glob import glob
import re

dir = input('講義回を2桁で入力してください : ')

# 文字列から数値だけを取り出す
def str2int(str):
    i = int( re.search('\d+',str).group() )
    return i

# ファイルの読み込み
fnms = glob("./" + dir + '/*.png')
# ソート
fnms.sort(key=str2int)
print(fnms)

# Presentationインスタンスの作成
ppt = Presentation()
# 幅
width = ppt.slide_width
# 高さ
height = ppt.slide_height

# レイアウト, 6番は白紙
blank_slide_layout = ppt.slide_layouts[6]


# ファイル毎にループ
for fnm in fnms:
    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(fnm, 0, 0)

    # 中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( ( height - pic.height ) / 2 )

# 名前をつけて保存
ppt.save(dir +'.pptx')